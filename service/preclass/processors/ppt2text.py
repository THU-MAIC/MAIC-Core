import os
import sys

from pymongo import MongoClient
from bson import ObjectId

from pptx import Presentation

from utils import get_channel, get_logger, now
from config import MONGO
from data.lecture import insert_file_snippet

import base64

def png_to_base64(png_file_path):
	"""
	Converts a PNG file to a Base64-encoded string.

	Parameters:
	png_file_path: The relative or absolute path to the PNG image file.

	Returns:
	A Base64-encoded string of the PNG image.
	"""
	try:
		# Open the PNG file in binary read mode
		with open(png_file_path, 'rb') as image_file:
			# Read the file's contents and encode it to Base64
			encoded_string = base64.b64encode(image_file.read())
		# Decode the Base64 bytes object to a string and return it
		return encoded_string.decode('utf-8')
	except FileNotFoundError:
		# Return an error message if the file is not found
		return "The specified file was not found. Please check the path is correct."
	except Exception as e:
		# Return a generic error message for any other exceptions
		return f"An error occurred during the conversion: {str(e)}"

def extract_text_from_ppt(
	ppt_path: str,
	png_path: str,
	lecture_id: ObjectId,
	) -> None:
	"""Extract text and images from PowerPoint slides and store them in the database.

	Args:
		ppt_path (str): Path to the PowerPoint file
		png_path (str): Directory path where PNG images will be stored
		lecture_id (ObjectId): MongoDB ObjectId of the lecture

	The function processes each slide to:
	- Extract text content from shapes
	- Convert slide to PNG image
	- Store both text and image in the database
	"""

	# Load the presentation
	presentation = Presentation(ppt_path)

	for slide_number, slide in enumerate(presentation.slides):
		slide_dir = os.path.join(png_path, str(slide_number + 1))
		slide_dir = slide_dir.replace("\\", "/")

		content = ""
		png = f"{slide_dir}.png"
		png_base64 = png_to_base64(png)

		# Extract text
		for shape in slide.shapes:
			if shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					content += paragraph.text + "\n"
				content += "\n"

		insert_file_snippet(
			idx=slide_number,
			content=content.strip(),
			png_base64=png_base64,
			lecture_id=lecture_id,
			file_type="pptx"
		)

class SERVICE:
	"""Service class for handling PowerPoint to text conversion tasks.

	Static Attributes:

			_collection (MongoClient=preclass.ppt2text): MongoDB collection for storing job information
				.. :noindex:
			_queue_name (str): RabbitMQ queue name for the service. Default to `preclass-ppt2text`
				.. :noindex:
			_logger: Logger instance for the service
				.. :noindex:
	"""

	_collection = MongoClient(
		MONGO.HOST,
		MONGO.PORT
		).preclass.ppt2text
	_queue_name = "preclass-ppt2text"

	_logger = get_logger(
		__name__=__name__,
		__file__=__file__,
	)

	@staticmethod
	def trigger(
			parent_service: str,
			lecture_id: ObjectId,
			parent_job_id: ObjectId
			) -> str:
		"""Trigger a new PPT to text conversion job.

		Args:
			parent_service (str): Name of the parent service
			lecture_id (ObjectId): MongoDB ObjectId of the lecture
			parent_job_id (ObjectId): MongoDB ObjectId of the parent job

		Returns:
			str: The job ID of the created conversion task
		"""
		connection, channel = get_channel(SERVICE._queue_name)
		
		SERVICE._logger.info("Pushing job to MONGO")
		
		job_id = SERVICE._collection.insert_one(
			dict(
				parent_service=parent_service,
				created_time = now(),
				lecture_id=lecture_id,
				parent_job_id=parent_job_id,
			)
		).inserted_id

		SERVICE._logger.info("Pushing job to RabbitMQ")
		channel.basic_publish(
			exchange="",
			routing_key=SERVICE._queue_name,
			body=str(job_id)
		)
		connection.close()
		
		SERVICE._logger.info("Job pushed to RabbitMQ")
		return job_id

	@staticmethod
	def callback(ch, method, properties, body):
		"""Process PowerPoint conversion jobs from the RabbitMQ queue.

		Args:
			ch: RabbitMQ channel
			method: RabbitMQ method frame
			properties: RabbitMQ properties
			body: Message body containing the job ID
		"""
		job_id = ObjectId(body.decode())
		job = SERVICE._collection.find_one(dict(_id=job_id))
		lecture_id, parent_service, parent_job_id = job["lecture_id"], job["parent_service"], job["parent_job_id"]
		SERVICE._logger.debug(f"Recieved PreClass PPT2TEXT Job - {lecture_id}")

		extract_text_from_ppt(
			ppt_path=f"buffer/{lecture_id}/seed_file.pptx",
			png_path=f"buffer/{lecture_id}/pngs",
			lecture_id=lecture_id
			)

		SERVICE._collection.update_one(
			dict(_id=job_id),
			{"$set": dict(
				completion_time=now()
			)}
		)
		SERVICE._logger.info(f"Conversion Complete For {lecture_id}")
		parent_connection, parent_channel = get_channel(parent_service)
		parent_channel.basic_publish(
			exchange="",
			routing_key=parent_service,
			body=str(parent_job_id)
		)
		parent_connection.close()
		ch.basic_ack(delivery_tag = method.delivery_tag)

	@staticmethod
	def launch_worker():
		"""Launch the worker to process PowerPoint conversion jobs.
		
		Starts consuming messages from the RabbitMQ queue and processes them.
		Can be terminated with CTRL+C.
		"""
		try:
			connection, channel = get_channel(SERVICE._queue_name)
			
			channel.basic_consume(
				queue=SERVICE._queue_name,
				on_message_callback=SERVICE.callback,
				auto_ack=False,
			)
			SERVICE._logger.info('Worker Launched. To exit press CTRL+C')
			channel.start_consuming()
		except KeyboardInterrupt:
			SERVICE._logger.warning('Shutting Off Worker')
			try:
				sys.exit(0)
			except SystemExit:
				os._exit(0)

if __name__ == "__main__":
	SERVICE._logger.warning("STARTING PRECLASS-PPTX2PDF SERVICE")
	SERVICE.launch_worker()